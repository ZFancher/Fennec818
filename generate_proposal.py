#!/usr/bin/env python3
"""
Generate ADS Plus proposal documents:
  1. ADS_Plus_Proposal.docx  — One-page Word proposal (shark tank style)
  2. ADS_Plus_Slides.pptx    — Two-slide visual presentation
"""

import os
BASE = '/Users/kyoti_m4/Desktop/ADS Plus Project'

# ──────────────────────────────────────────────────────────────────────────────
#  WORD DOCUMENT
# ──────────────────────────────────────────────────────────────────────────────

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

NAVY_HEX    = '1B3260'
GOLD_HEX    = 'C9A023'
GREEN_HEX   = '4B5320'
RED_HEX     = 'A01E1E'

NAVY        = (27,  50,  96)
GOLD        = (201, 160,  35)
DKGREEN     = (75,  83,  32)
WHITE       = (255, 255, 255)
DGRAY       = (51,  51,  51)
MGRAY       = (120, 120, 120)
RED         = (160,  30,  30)


def cell_shade(cell, hex_color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), hex_color)
    tcPr.append(shd)


def cell_no_border(cell):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcBorders = OxmlElement('w:tcBorders')
    for side in ['top', 'bottom', 'left', 'right']:
        b = OxmlElement(f'w:{side}')
        b.set(qn('w:val'), 'nil')
        tcBorders.append(b)
    tcPr.append(tcBorders)


def table_no_border(table):
    tbl = table._tbl
    tblPr = tbl.tblPr
    if tblPr is None:
        tblPr = OxmlElement('w:tblPr')
        tbl.insert(0, tblPr)
    tblBorders = OxmlElement('w:tblBorders')
    for name in ['top', 'bottom', 'left', 'right', 'insideH', 'insideV']:
        b = OxmlElement(f'w:{name}')
        b.set(qn('w:val'), 'none')
        tblBorders.append(b)
    tblPr.append(tblBorders)


def run(para, text, size=9.5, bold=False, color=None, italic=False):
    r = para.add_run(text)
    r.font.size = Pt(size)
    r.bold = bold
    r.italic = italic
    if color:
        r.font.color.rgb = RGBColor(*color)
    return r


def fmt(p, sb=0, sa=2, ls=None, li=None, fi=None, align=None):
    pf = p.paragraph_format
    pf.space_before = Pt(sb)
    pf.space_after  = Pt(sa)
    if ls is not None: pf.line_spacing = Pt(ls)
    if li is not None: pf.left_indent  = Inches(li)
    if fi is not None: pf.first_line_indent = Inches(fi)
    if align: p.alignment = align


def section_hdr(doc, text, hex_color):
    t = doc.add_table(rows=1, cols=1)
    table_no_border(t)
    c = t.cell(0, 0)
    cell_shade(c, hex_color)
    cell_no_border(c)
    p = c.paragraphs[0]
    fmt(p, sb=2, sa=2, li=0.1)
    run(p, text, size=9, bold=True, color=WHITE)


def two_col_table(doc):
    t = doc.add_table(rows=1, cols=2)
    table_no_border(t)
    lc, rc = t.cell(0, 0), t.cell(0, 1)
    cell_no_border(lc); cell_no_border(rc)
    lc.width = Inches(3.5); rc.width = Inches(3.5)
    return lc, rc


def add_bullet(cell, is_first, symbol, label, body, sym_color, label_color=DGRAY, size=9, ls=12, indent=0.12):
    p = cell.paragraphs[0] if is_first else cell.add_paragraph()
    fmt(p, sb=2, sa=1, ls=ls, li=indent, fi=-indent)
    run(p, f'{symbol} ', size=size, color=sym_color, bold=True)
    if label:
        run(p, f'{label} ', size=size, bold=True, color=label_color)
    run(p, body, size=size)
    return p


def create_docx(path):
    doc = Document()

    # Default style tweaks
    doc.styles['Normal'].font.name = 'Calibri'
    doc.styles['Normal'].paragraph_format.space_before = Pt(0)
    doc.styles['Normal'].paragraph_format.space_after  = Pt(0)

    for sec in doc.sections:
        sec.page_height    = Inches(11)
        sec.page_width     = Inches(8.5)
        sec.top_margin     = Inches(0.5)
        sec.bottom_margin  = Inches(0.5)
        sec.left_margin    = Inches(0.75)
        sec.right_margin   = Inches(0.75)

    # ── HEADER BANNER ──────────────────────────────────────────────────────────
    ht = doc.add_table(rows=1, cols=2)
    table_no_border(ht)
    lc, rc = ht.cell(0, 0), ht.cell(0, 1)
    lc.width = Inches(5.2); rc.width = Inches(1.8)
    cell_shade(lc, NAVY_HEX); cell_shade(rc, NAVY_HEX)
    cell_no_border(lc); cell_no_border(rc)

    lp = lc.paragraphs[0]
    fmt(lp, sb=5, sa=0, li=0.12)
    run(lp, 'ADS PLUS', size=20, bold=True, color=WHITE)

    lp2 = lc.add_paragraph()
    fmt(lp2, sb=1, sa=4, li=0.12)
    run(lp2, 'Streamlining Arid West Wetland Delineation Through Mobile AI', size=9, italic=True, color=GOLD)

    rp = rc.paragraphs[0]
    fmt(rp, sb=5, sa=1, align=WD_ALIGN_PARAGRAPH.RIGHT)
    run(rp, 'USACE Regulatory Program', size=7.5, color=WHITE)
    rp2 = rc.add_paragraph()
    fmt(rp2, sb=0, sa=1, align=WD_ALIGN_PARAGRAPH.RIGHT)
    run(rp2, 'Arid West Region', size=7.5, color=WHITE)
    rp3 = rc.add_paragraph()
    fmt(rp3, sb=0, sa=4, align=WD_ALIGN_PARAGRAPH.RIGHT)
    run(rp3, 'June 2026', size=7.5, bold=True, color=GOLD)

    # Gold rule
    sep = doc.add_paragraph()
    fmt(sep, sb=0, sa=3)
    pPr = sep._p.get_or_add_pPr()
    pBdr = OxmlElement('w:pBdr')
    bot = OxmlElement('w:bottom')
    bot.set(qn('w:val'), 'single'); bot.set(qn('w:sz'), '8')
    bot.set(qn('w:space'), '0'); bot.set(qn('w:color'), GOLD_HEX)
    pBdr.append(bot); pPr.append(pBdr)

    # ── THE PROBLEM ────────────────────────────────────────────────────────────
    section_hdr(doc, 'THE PROBLEM — Administrative Overhead in a Multi-Tool Regulatory Workflow', NAVY_HEX)

    lc, rc = two_col_table(doc)

    # Left: narrative + bullets
    p0 = lc.paragraphs[0]
    fmt(p0, sb=2, sa=2, ls=12, li=0.08)
    run(p0, 'USACE regulatory staff conducting wetland delineations in the Arid West currently navigate '
        'three separate tools across two environments: ', size=9)
    run(p0, 'Survey 123 ', size=9, bold=True)
    run(p0, 'for field data entry, the ')
    run(p0, 'National Regulatory Viewer (NRV) ', size=9, bold=True)
    run(p0, 'for transfer, and the ')
    run(p0, 'ADS Excel workbook ', size=9, bold=True)
    run(p0, 'for inference and form output. This fragmented pipeline creates compounding inefficiencies '
        'that directly impact mission delivery and data quality.')

    prob_bullets = [
        ('Post-field re-entry burden:', 'Data from Survey 123 must be manually transferred through the NRV '
         'into the ADS before any inference is available — requiring return to the office after every field session.'),
        ('Silent data contradictions:', 'Indicator selections made in Survey 123 can conflict with ADS '
         'auto-inference on the same data. No flag, no alert, no reconciliation step exists on the final ENG Form 6116-1.'),
        ('Formula fragility:', 'The ADS Excel formula network can be corrupted by routine operations '
         '(copy-paste, data entry) — producing incorrect determinations with no user-visible warning.'),
        ('Zero field guidance:', 'Neither tool explains why an indicator fires or fails. Delineators work '
         'without a real-time QC or reference system in the field.'),
        ('No camera integration:', 'Soil color requires a physical Munsell chart, carried separately, '
         'read subjectively under variable lighting — a significant source of inter-observer error.'),
    ]
    for i, (lbl, body) in enumerate(prob_bullets):
        add_bullet(lc, False, '▸', lbl, body, sym_color=NAVY, label_color=NAVY, size=9)

    # Right: workflow
    rp0 = rc.paragraphs[0]
    fmt(rp0, sb=2, sa=2, li=0.08)
    run(rp0, 'CURRENT WORKFLOW', size=8.5, bold=True, color=NAVY)

    steps = [
        ('FIELD',   'Survey 123 (Esri)',              'Mobile entry — no inference, no QC, no guidance'),
        ('OFFICE',  'NRV Upload',                     'Manual post-field data transfer — connectivity required'),
        ('OFFICE',  'ADS Excel Workbook',             'Inference & PDF — laptop required, formula-fragile'),
        ('OUTPUT',  'ENG Form 6116-1 via NRV',        'Risk of Survey 123 / ADS data contradiction on final form'),
    ]
    for i, (tag, name, note) in enumerate(steps):
        sp = rc.add_paragraph()
        fmt(sp, sb=3, sa=0, ls=12, li=0.08)
        run(sp, f'{tag}  ', size=7, color=MGRAY, bold=True)
        run(sp, name, size=9.5, bold=True, color=NAVY)
        sn = rc.add_paragraph()
        fmt(sn, sb=0, sa=0, ls=11, li=0.22)
        run(sn, f'↳ {note}', size=8, italic=True, color=RED)
        if i < 3:
            arr = rc.add_paragraph()
            fmt(arr, sb=0, sa=0, li=0.18)
            run(arr, '↓', size=9, color=MGRAY)

    # ── THE AI-POWERED SOLUTION ────────────────────────────────────────────────
    section_hdr(doc, 'THE AI-POWERED SOLUTION — ADS Plus', GREEN_HEX)

    ip = doc.add_paragraph()
    fmt(ip, sb=2, sa=2, ls=12, li=0.08)
    run(ip, 'ADS Plus ', size=9, bold=True)
    run(ip, 'is a browser-based, offline-capable application that consolidates the entire three-factor '
        'wetland delineation workflow — vegetation, hydric soils, and hydrology — into a single HTML file '
        'running on any smartphone, tablet, or laptop with no installation, server, or internet connection required.',
        size=9)

    lc2, rc2 = two_col_table(doc)

    sol_left = [
        ('AI Camera — Color Correction Matrix (CCM):',
         'A printed reference card with known Lab values enables real-time least-squares regression to derive '
         'a per-frame 3×3 Color Correction Matrix, correcting camera spectral response for any field lighting. '
         'Replaces subjective physical Munsell chart reading.'),
        ('AI Camera — Automated Redox Detection (SoilAnalyzeCapture):',
         'K-means++ clustering (K=6) partitions the full color distribution of a ped face. Clusters are matched '
         'to Munsell chips via CIE76 ΔE. Concentrations classified by USDA Table A1 contrast; depletions by '
         'absolute Munsell chroma. Outputs matrix color, feature type, Munsell notation, and % area from a '
         'single photograph — a capability with no equivalent in any existing USACE field tool.'),
        ('Automated Indicator Inference with Reason Strings:',
         'Evaluates all applicable hydric soil and hydrology indicators from entered data, returning '
         'plain-language explanations of exactly which horizon at which depth with which measurements '
         'triggered each indicator — supporting both documentation and field training.'),
        ('Data Consistency Enforcement:',
         'Real-time warnings surface when manually selected indicators conflict with entered data, '
         'LRR restrictions, above-layer chroma rules, or problematic indicator conditions — '
         'preventing the silent contradictions inherent in the multi-tool pipeline.'),
    ]
    for i, (lbl, body) in enumerate(sol_left):
        add_bullet(lc2, i == 0, '▸', lbl, body, sym_color=DKGREEN, label_color=DKGREEN, size=9)

    sol_right = [
        ('In-App ENG Form 6116-1 PDF Generation:',
         'A filled, print-ready data sheet generated directly within the application. Eliminates the NRV '
         'workflow step for form output entirely.'),
        ('Complete Offline Operation:',
         'Single HTML file. No installation, no App Store, no server, no internet connection after '
         'initial download. All data stored locally via browser localStorage.'),
        ('Full JSON State Export:',
         'Every horizon, selection, inference result, GPS coordinate, and remark exportable to a '
         'structured JSON archive — providing a complete, auditable, peer-reviewable data record.'),
        ('Single-Application Workflow:',
         'Replaces Survey 123 → NRV → ADS → NRV entirely. Field data entry through PDF output — '
         'one app, one session, one device, zero post-field transfer.'),
    ]
    for i, (lbl, body) in enumerate(sol_right):
        add_bullet(rc2, i == 0, '▸', lbl, body, sym_color=DKGREEN, label_color=DKGREEN, size=9)

    # ── EXPECTED IMPACT / IMPLEMENTATION ──────────────────────────────────────
    section_hdr(doc, 'EXPECTED IMPACT  &  IMPLEMENTATION REQUIREMENTS', '2D4A1A')

    lc3, rc3 = two_col_table(doc)

    impact_hdr = lc3.paragraphs[0]
    fmt(impact_hdr, sb=2, sa=2, li=0.08)
    run(impact_hdr, 'EXPECTED IMPACT', size=8.5, bold=True, color=DKGREEN)

    impact_items = [
        'Eliminates multi-step post-field pipeline; estimated 30–60 min saved per delineation point',
        'Removes structural risk of cross-tool data contradictions on the ENG Form 6116-1',
        'Increases technical accuracy: implements current NRCS Field Indicators 9.3 depth thresholds for F3 '
        'and F8, correcting known formula holdovers in the ADS',
        'Context-aware inference explanations accelerate delineator training',
        'No hardware cost beyond a government-issue smartphone or tablet',
    ]
    for item in impact_items:
        add_bullet(lc3, False, '•', None, item, sym_color=DKGREEN, size=9, ls=12, indent=0.15)

    impl_hdr = rc3.paragraphs[0]
    fmt(impl_hdr, sb=2, sa=2, li=0.08)
    run(impl_hdr, 'IMPLEMENTATION REQUIREMENTS', size=8.5, bold=True, color=NAVY)

    impl_items = [
        ('Platform:', 'Any modern smartphone, tablet, or laptop (iOS, Android, Windows, macOS)'),
        ('Distribution:', 'Single HTML file — no App Store, no installation, no configuration'),
        ('Network:', 'None required for operation; one-time file transfer for initial deployment'),
        ('Camera (AI modules):', 'Standard rear-facing smartphone camera with autofocus'),
        ('Reference Card:', 'Printed on a standard color printer; template included in app'),
        ('Approval Authority:', 'Formal adoption requires USACE Regulatory Program and HQ USACE coordination'),
    ]
    for lbl, body in impl_items:
        p = rc3.add_paragraph()
        fmt(p, sb=1, sa=1, ls=12, li=0.15, fi=-0.15)
        run(p, '• ', size=9, color=NAVY, bold=True)
        run(p, f'{lbl} ', size=9, bold=True)
        run(p, body, size=9)

    # ── CYBERSECURITY / COMPLIANCE ─────────────────────────────────────────────
    section_hdr(doc, 'CYBERSECURITY, RECORDS MANAGEMENT & POLICY COMPLIANCE', RED_HEX)

    comp_items = [
        ('Cybersecurity (AR 25-2):', 'All data remains on-device via browser localStorage — '
         'no data transmitted to any server, cloud service, or external system. No ATO required for '
         'local browser execution.'),
        ('Records Management (AR 25-400-2):', 'Full JSON export provides a timestamped, auditable '
         'record of all entered data, selections, and calculated results, consistent with Army records '
         'management requirements.'),
        ('Privacy (DoD Privacy Program):', 'No PII collected or transmitted. GPS coordinates stored '
         'locally; included only in user-initiated exports under delineator control.'),
        ('AI Governance (DoD AI Strategy / EO 13960):', 'AI capabilities — CCM color correction, '
         'k-means clustering, indicator inference — are transparent, explainable, and human-supervised. '
         'The delineator reviews all results and retains final determination authority at all times.'),
        ('Authoritative Basis:', '1987 USACE Wetlands Delineation Manual; Arid West Regional Supplement '
         '(2008); NRCS Field Indicators of Hydric Soils Version 9.3 (Feb 2026 errata); ENG Form 6116-1 '
         '(Sep 2024).'),
    ]
    for lbl, body in comp_items:
        p = doc.add_paragraph()
        fmt(p, sb=2, sa=1, ls=12, li=0.12, fi=-0.12)
        run(p, '▸ ', size=9, color=RED, bold=True)
        run(p, f'{lbl} ', size=9, bold=True)
        run(p, body, size=9)

    # ── FOOTER ─────────────────────────────────────────────────────────────────
    fp = doc.add_paragraph()
    fmt(fp, sb=4, sa=0, align=WD_ALIGN_PARAGRAPH.CENTER)
    pPr2 = fp._p.get_or_add_pPr()
    pBdr2 = OxmlElement('w:pBdr')
    top = OxmlElement('w:top')
    top.set(qn('w:val'), 'single'); top.set(qn('w:sz'), '6')
    top.set(qn('w:space'), '1'); top.set(qn('w:color'), NAVY_HEX)
    pBdr2.append(top); pPr2.append(pBdr2)
    run(fp, 'ADS Plus v15  |  USACE Arid West Regulatory Program  |  Innovation Proposal  |  June 2026  |  UNCLASSIFIED',
        size=7.5, color=MGRAY)

    doc.save(path)
    print(f'Word document saved → {path}')


# ──────────────────────────────────────────────────────────────────────────────
#  POWERPOINT
# ──────────────────────────────────────────────────────────────────────────────

from pptx import Presentation
from pptx.util import Inches as Pi, Pt as PPt
from pptx.dml.color import RGBColor as Pc
from pptx.enum.text import PP_ALIGN

# Color tuples for pptx (use Pc(*tuple) to get RGBColor)
PNAVY   = (27,  50,  95)
PGOLD   = (201, 160,  35)
PGREEN  = (75,  83,  32)
PWHITE  = (255, 255, 255)
PLGRAY  = (240, 242, 245)
PDGRAY  = (51,  51,  51)
PAMBER  = (200, 110,  25)
PTEAL   = (0,   130, 120)
PSOIL   = (118,  92,  58)
PSOIL2  = (100,  76,  44)
PMGRAY  = (148, 143, 138)
PRED    = (160,  35,  35)
PSILVER = (190, 190, 195)
PLGREEN = (100, 180,  60)
PCREAM  = (240, 235, 225)


def prect(slide, x, y, w, h, fill, line=None, lw=None):
    shp = slide.shapes.add_shape(1, Pi(x), Pi(y), Pi(w), Pi(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = Pc(*fill)
    if line:
        shp.line.color.rgb = Pc(*line)
        if lw: shp.line.width = PPt(lw)
    else:
        shp.line.fill.background()
    return shp


def prrect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(5, Pi(x), Pi(y), Pi(w), Pi(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = Pc(*fill)
    if line: shp.line.color.rgb = Pc(*line)
    else: shp.line.fill.background()
    return shp


def ptxt(slide, text, x, y, w, h, size=12, color=PWHITE, bold=False,
         align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Pi(x), Pi(y), Pi(w), Pi(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size  = PPt(size)
        r.font.bold  = bold
        r.font.italic = italic
        r.font.color.rgb = Pc(*color)
    return tb


def pgold_rule(slide, x, y, w):
    prect(slide, x, y, w, 0.03, PGOLD)


def create_pptx(path):
    prs = Presentation()
    prs.slide_width  = Pi(13.33)
    prs.slide_height = Pi(7.5)
    blank = prs.slide_layouts[6]

    # ═══════════════════════════════════════════════════════════════════════════
    #  SLIDE 1 — THE PROBLEM
    # ═══════════════════════════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(blank)

    # Background
    prect(s1, 0, 0, 13.33, 7.5, PLGRAY)

    # Header band
    prect(s1, 0, 0, 13.33, 1.05, PNAVY)
    ptxt(s1, 'ADS PLUS', 0.2, 0.07, 4.5, 0.5, size=28, bold=True, color=PWHITE)
    ptxt(s1, 'Innovation Proposal  |  USACE Regulatory Program  |  Arid West Region  |  June 2026',
         0.22, 0.58, 9.0, 0.35, size=9, color=PGOLD, italic=True)
    ptxt(s1, 'UNCLASSIFIED', 11.6, 0.1, 1.5, 0.3, size=8, color=PWHITE, bold=True, align=PP_ALIGN.RIGHT)
    ptxt(s1, '1 / 2', 12.85, 0.68, 0.38, 0.25, size=8, color=PSILVER, align=PP_ALIGN.RIGHT)

    # Slide title area
    ptxt(s1, 'THE ADMINISTRATIVE BURDEN', 0.2, 1.12, 9.5, 0.48, size=20, bold=True, color=PNAVY)
    ptxt(s1, 'Why the current digital delineation workflow creates mission-critical obstacles',
         0.22, 1.62, 9.5, 0.3, size=10.5, color=PDGRAY, italic=True)
    pgold_rule(s1, 0.2, 1.97, 12.9)

    # ── LEFT PANEL: Current Workflow Diagram ───────────────────────────────────
    prect(s1, 0.2, 2.07, 5.5, 5.08, PWHITE, line=PNAVY, lw=0.5)
    ptxt(s1, 'CURRENT THREE-TOOL WORKFLOW', 0.35, 2.18, 5.2, 0.28, size=8.5, bold=True,
         color=PNAVY, align=PP_ALIGN.CENTER)
    ptxt(s1, '(Survey 123  →  NRV  →  ADS Excel  →  ENG 6116-1)',
         0.35, 2.44, 5.2, 0.22, size=8, color=PMGRAY, align=PP_ALIGN.CENTER, italic=True)

    step_defs = [
        (PGREEN,          'FIELD',   'Survey 123 (Esri Mobile)',
         'No auto-inference  •  No QC checks  •  No guidance  •  Manual only'),
        ((70, 80, 105),   'OFFICE',  'National Regulatory Viewer (NRV) Upload',
         'Internet required  •  Post-field only  •  Manual data transfer step'),
        ((70, 80, 105),   'OFFICE',  'ADS Excel Workbook',
         'Laptop required  •  Formula-fragile  •  No field optimization  •  No camera'),
        (PNAVY,           'OUTPUT',  'ENG Form 6116-1 Output',
         'Risk of Survey 123 / ADS data contradictions  •  No in-field output possible'),
    ]

    by = 2.75
    for col, tag, name, note in step_defs:
        prrect(s1, 0.3, by, 0.72, 0.34, col)
        ptxt(s1, tag, 0.3, by + 0.04, 0.72, 0.28, size=7.5, bold=True, align=PP_ALIGN.CENTER)
        prect(s1, 1.1, by, 4.42, 0.34, col)
        ptxt(s1, name, 1.15, by + 0.04, 4.35, 0.28, size=9.5, bold=True)
        ptxt(s1, note, 1.15, by + 0.38, 4.4, 0.24, size=8, color=PRED, italic=True)
        by += 1.0
        if by < 5.75:
            prect(s1, 2.38, by - 0.55, 0.09, 0.35, PNAVY)
            prect(s1, 2.3,  by - 0.22, 0.25, 0.09, PNAVY)

    # ── CENTER PANEL: Pain Points ─────────────────────────────────────────────
    ptxt(s1, '⚠  KEY PAIN POINTS', 5.9, 2.07, 3.7, 0.35, size=10, bold=True, color=PNAVY)

    pains = [
        ('30–60 min',        'post-field re-entry per delineation point — before any inference is available'),
        ('Silent data\nconflicts', 'Survey 123 manual selections can contradict ADS auto-inference on the same data with no alert on the final 6116-1'),
        ('Formula\nfragility', 'ADS Excel formula network silently corrupted by routine copy-paste or data entry — wrong output, no warning'),
        ('Zero field\nguidance', 'No tool explains why an indicator fires or fails — no real-time QC or reference system in remote terrain'),
        ('Physical\nMunsell charts', 'Soil color requires a separate, subjectively-read chart — no camera tool exists in any current USACE system'),
        ('Connectivity\ndependency', 'Survey 123 requires internet — creating data gaps in remote Arid West terrain where cell coverage is absent'),
    ]

    py = 2.52
    for stat, desc in pains:
        prect(s1, 5.9, py, 3.65, 0.62, PNAVY)
        ptxt(s1, stat,  5.95, py + 0.03, 3.55, 0.28, size=10.5, bold=True, color=PGOLD)
        ptxt(s1, desc,  5.95, py + 0.32, 3.55, 0.27, size=8, color=PWHITE, wrap=True)
        py += 0.72

    # ── RIGHT PANEL: Mission Impact ────────────────────────────────────────────
    prect(s1, 9.65, 2.07, 3.5, 5.08, PNAVY)
    ptxt(s1, 'MISSION IMPACT', 9.8, 2.16, 3.2, 0.3, size=10, bold=True, color=PGOLD)

    impact_text = (
        'This fragmented workflow consumes significant staff '
        'time, introduces structural data quality risks, and '
        'limits USACE regulatory staff from producing timely, '
        'defensible wetland determinations in the field.\n\n'
        'The multi-tool pipeline produces regulatory outputs '
        'through a process with no integrated QC, no field '
        'verification, and no transparent chain of reasoning '
        'from field data to final determination.\n\n'
        'Every sample point in the Arid West carries this '
        'overhead. Across a delineation season of dozens of '
        'points, the cumulative burden is substantial and '
        'error risk compounds with each transfer step.\n\n'
        'The current toolset was not designed for the field '
        'delineator. ADS Plus was.'
    )
    ptxt(s1, impact_text, 9.78, 2.55, 3.25, 4.35, size=9.5, color=PWHITE, wrap=True)

    # Footer
    prect(s1, 0, 7.18, 13.33, 0.32, PNAVY)
    ptxt(s1, 'ADS Plus  |  A Next-Generation Wetland Delineation Support Tool  |  USACE Regulatory Program  |  UNCLASSIFIED',
         0.2, 7.2, 13.0, 0.26, size=8, color=PSILVER, align=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════════════════════
    #  SLIDE 2 — THE SOLUTION
    # ═══════════════════════════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(blank)

    # Background
    prect(s2, 0, 0, 13.33, 7.5, PLGRAY)

    # Header band — Army green
    prect(s2, 0, 0, 13.33, 1.05, PGREEN)
    ptxt(s2, 'ADS PLUS', 0.2, 0.07, 4.5, 0.5, size=28, bold=True, color=PWHITE)
    ptxt(s2, 'Innovation Proposal  |  USACE Regulatory Program  |  Arid West Region  |  June 2026',
         0.22, 0.58, 9.0, 0.35, size=9, color=PGOLD, italic=True)
    ptxt(s2, 'UNCLASSIFIED', 11.6, 0.1, 1.5, 0.3, size=8, color=PWHITE, bold=True, align=PP_ALIGN.RIGHT)
    ptxt(s2, '2 / 2', 12.85, 0.68, 0.38, 0.25, size=8, color=PSILVER, align=PP_ALIGN.RIGHT)

    # Slide title
    ptxt(s2, 'ONE APP.  ONE SESSION.  ONE FORM.', 0.2, 1.12, 10.0, 0.48, size=20, bold=True, color=PGREEN)
    ptxt(s2, 'ADS Plus consolidates the entire delineation workflow — vegetation, soils, hydrology, inference, and output — in a single offline-capable app',
         0.22, 1.62, 10.0, 0.3, size=10.5, color=PDGRAY, italic=True)
    pgold_rule(s2, 0.2, 1.97, 12.9)

    # ── LEFT PANEL: Capabilities ──────────────────────────────────────────────
    prect(s2, 0.2, 2.07, 6.6, 5.08, PWHITE, line=PGREEN, lw=0.5)
    ptxt(s2, 'CAPABILITIES', 0.35, 2.18, 6.3, 0.28, size=9, bold=True, color=PGREEN)

    caps = [
        ('AI Camera: Color Correction Matrix (CCM)',
         'Printed reference card with known Lab values → per-frame least-squares regression → '
         '3×3 CCM correcting camera spectral response for any field lighting. Enables accurate '
         'Munsell color matching in sun, shade, or mixed light with no user action. '
         'Replaces subjective physical chart reading.'),
        ('AI Camera: Automated Redox Detection (SoilAnalyzeCapture)',
         'K-means++ clustering (K=6) segments full color distribution of a soil ped face. '
         'Each cluster matched to Munsell chip via CIE76 ΔE. Concentrations classified by '
         'USDA Table A1 contrast; depletions by Munsell chroma. Returns matrix color, '
         'feature type, Munsell notation, and % area — from a single photograph. '
         'No equivalent exists in any current USACE field tool.'),
        ('Automated Indicator Inference with Reason Strings',
         'Evaluates all applicable hydric soil and hydrology indicators from entered data. '
         'Returns plain-language explanations: which horizon at which depth with which '
         'measured properties triggered each indicator. Supports documentation and '
         'accelerates delineator training.'),
        ('Data Consistency Enforcement & In-App PDF',
         'Real-time warnings when manually selected indicators conflict with entered data, '
         'LRR rules, or chroma requirements. In-app ENG Form 6116-1 PDF generated directly — '
         'no NRV upload, no post-field transfer, no second device. '
         'Complete offline operation via single HTML file.'),
    ]

    cy = 2.5
    for title, body in caps:
        prect(s2, 0.25, cy, 0.06, 0.72, PGREEN)
        ptxt(s2, title, 0.38, cy,       6.35, 0.26, size=9.5, bold=True, color=PGREEN)
        ptxt(s2, body,  0.38, cy + 0.27, 6.35, 0.5,  size=8.5, color=PDGRAY, wrap=True)
        cy += 1.24

    # ── RIGHT PANEL: AI Visual Modules ────────────────────────────────────────

    # ── Visual Panel 1: Color Calibration Reference Card ──────────────────────
    prect(s2, 6.98, 2.07, 6.18, 2.58, PNAVY)
    ptxt(s2, 'AI MODULE 1 — COLOR CALIBRATION REFERENCE CARD',
         7.1, 2.14, 5.95, 0.28, size=8.5, bold=True, color=PGOLD)
    ptxt(s2, 'Known Lab-value patches  •  CCM solved per frame via least-squares  •  Adapts to any field lighting',
         7.1, 2.4, 5.95, 0.22, size=7.5, color=PSILVER, italic=True)

    # Draw the printed reference card
    cx, cy2, cw, ch = 7.1, 2.68, 3.9, 1.78
    prect(s2, cx, cy2, cw, ch, PCREAM, line=(200, 195, 190), lw=1.5)

    # Color patch grid: 4 rows × 6 cols (representing calibration patches)
    patch_grid = [
        # Row 1: warm reds to oranges (concentration zone spectrum)
        [(155, 55, 50), (180, 82, 46), (196, 115, 50), (202, 140, 55), (192, 160, 76), (178, 155, 112)],
        # Row 2: soil browns / yellows (matrix color spectrum)
        [(120, 85, 50), (140, 104, 64), (155, 120, 75), (165, 140, 90), (170, 150, 110), (175, 165, 135)],
        # Row 3: grays / neutral (depletion spectrum)
        [(80, 75, 72), (110, 105, 100), (142, 138, 133), (170, 165, 160), (200, 196, 192), (225, 220, 216)],
        # Row 4: dark reference → near white
        [(28, 28, 28), (68, 63, 58), (108, 96, 84), (142, 130, 118), (188, 180, 170), (240, 237, 230)],
    ]
    pad, gap = 0.08, 0.025
    pw2 = (cw - 2 * pad) / 6 - gap
    ph2 = (ch - 2 * pad) / 4 - gap

    for ri, row in enumerate(patch_grid):
        for ci, pc2 in enumerate(row):
            prect(s2,
                  cx + pad + ci * (pw2 + gap),
                  cy2 + pad + ri * (ph2 + gap),
                  pw2, ph2, pc2)

    # ArUco fiducial markers at each corner (black outer, white middle, black center)
    mk = 0.15
    for mx2, my2 in [(cx + 0.02, cy2 + 0.02),
                     (cx + cw - mk - 0.02, cy2 + 0.02),
                     (cx + 0.02, cy2 + ch - mk - 0.02),
                     (cx + cw - mk - 0.02, cy2 + ch - mk - 0.02)]:
        prect(s2, mx2, my2, mk, mk, (10, 10, 10))
        prect(s2, mx2 + 0.033, my2 + 0.033, mk - 0.066, mk - 0.066, PWHITE)
        prect(s2, mx2 + 0.053, my2 + 0.053, mk - 0.106, mk - 0.106, (10, 10, 10))

    # Status readout next to card
    prect(s2, 11.12, 2.72, 1.9, 0.36, (30, 110, 55))
    ptxt(s2, '✓  CCM DETECTED',  11.17, 2.76, 1.8, 0.28, size=9, bold=True, color=PWHITE)
    ptxt(s2, 'ΔE confidence: 2.1',    11.17, 3.12, 1.8, 0.22, size=8, color=PSILVER)
    ptxt(s2, 'Sharpness check: PASS', 11.17, 3.32, 1.8, 0.22, size=8, color=PSILVER)
    ptxt(s2, 'Lighting: FIELD CORRECTED', 11.17, 3.52, 1.8, 0.22, size=8, color=PGOLD)

    # ── Visual Panel 2: Redox Detection Overlay ────────────────────────────────
    prect(s2, 6.98, 4.72, 6.18, 2.43, PNAVY)
    ptxt(s2, 'AI MODULE 2 — AUTOMATED REDOX DETECTION OVERLAY',
         7.1, 4.79, 5.95, 0.28, size=8.5, bold=True, color=PGOLD)
    ptxt(s2, 'K-means++ clustering  •  USDA Table A1 contrast classification  •  Real-time slider exploration',
         7.1, 5.05, 5.95, 0.22, size=7.5, color=PSILVER, italic=True)

    # Ped face (matrix color ~10YR 4/2 grayish brown)
    px2, py3, pw3, ph3 = 7.12, 5.32, 3.15, 1.72
    prect(s2, px2, py3, pw3, ph3, PSOIL)

    # Redox concentration patches (amber — ~7.5YR 4/6)
    concs = [
        (0.12, 0.08, 0.52, 0.24), (0.75, 0.14, 0.43, 0.28), (0.28, 0.52, 0.38, 0.22),
        (1.28, 0.04, 0.36, 0.20), (1.75, 0.44, 0.48, 0.28), (2.08, 0.08, 0.34, 0.18),
        (0.58, 0.82, 0.42, 0.22), (1.48, 0.72, 0.38, 0.24), (0.92, 1.25, 0.3, 0.2),
    ]
    for rx, ry, rw2, rh2 in concs:
        prect(s2, px2 + rx, py3 + ry, rw2, rh2, PAMBER)

    # Depletion patches (teal ~10YR 5/1)
    depls = [(2.5, 0.28, 0.52, 0.36), (2.05, 1.2, 0.62, 0.3)]
    for rx, ry, rw2, rh2 in depls:
        prect(s2, px2 + rx, py3 + ry, rw2, rh2, PTEAL)

    # Selection polygon outline (yellow dashed simulation — thin rect with no fill)
    sel = s2.shapes.add_shape(1, Pi(px2 + 0.06), Pi(py3 + 0.04), Pi(pw3 - 0.12), Pi(ph3 - 0.08))
    sel.fill.background()
    sel.line.color.rgb = Pc(255, 220, 40)
    sel.line.width = PPt(1.5)

    # Results panel
    ptxt(s2, 'RESULTS', 10.38, 5.32, 2.72, 0.26, size=8.5, bold=True, color=PGOLD)

    prect(s2, 10.38, 5.6,  2.7, 0.22, PSOIL)
    ptxt(s2, 'Matrix  10YR 4/2   58%', 10.42, 5.62, 2.62, 0.2, size=8.5, bold=True, color=PWHITE)

    prect(s2, 10.38, 5.86, 0.18, 0.18, PAMBER)
    ptxt(s2, 'Conc  7.5YR 4/6   22%  Distinct', 10.6, 5.85, 2.45, 0.22, size=8.5, bold=True, color=PWHITE)

    prect(s2, 10.38, 6.12, 0.18, 0.18, PTEAL)
    ptxt(s2, 'Depl  10YR 5/1   9%', 10.6, 6.11, 2.45, 0.22, size=8.5, bold=True, color=PWHITE)

    prect(s2, 10.38, 6.42, 2.7, 0.3, (30, 110, 55))
    ptxt(s2, '✓  F3 Depleted Matrix AUTO-INFERRED', 10.42, 6.46, 2.62, 0.24,
         size=8, bold=True, color=PWHITE)

    ptxt(s2, 'Chroma ≤2  •  ≥6 in.  •  above-layer check passed',
         10.42, 6.76, 2.62, 0.22, size=7.5, color=PSILVER, italic=True)

    # Footer outcome strip
    prect(s2, 0, 7.18, 13.33, 0.32, PGREEN)
    ptxt(s2, '✓ Eliminates multi-tool pipeline    '
         '✓ AI-powered camera soil color analysis    '
         '✓ Real-time indicator inference with explanations    '
         '✓ In-app ENG Form 6116-1 PDF    '
         '✓ No infrastructure required    '
         '✓ Fully offline',
         0.2, 7.2, 13.0, 0.26, size=8.5, bold=True, color=PWHITE, align=PP_ALIGN.CENTER)

    prs.save(path)
    print(f'PowerPoint saved → {path}')


# ──────────────────────────────────────────────────────────────────────────────
if __name__ == '__main__':
    docx_path = os.path.join(BASE, 'ADS_Plus_Proposal.docx')
    pptx_path = os.path.join(BASE, 'ADS_Plus_Slides.pptx')
    create_docx(docx_path)
    create_pptx(pptx_path)
    print('Done.')
